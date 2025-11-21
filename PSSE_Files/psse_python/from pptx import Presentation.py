from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12.0), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00,0x33,0x66)
    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.7))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(0x33,0x33,0x33)

def add_bullets_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(0.7)
    top = Inches(1.6)
    width = Inches(12)
    height = Inches(5.0)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Pt(6)
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
            p.level = 0
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x11,0x11,0x11)

def add_flowchart_slide(title, steps, start_x=1.0, start_y=1.0, box_w=3.5, box_h=0.7, gap=0.6):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.title.text = title
    # add boxes vertically with connectors
    left = Inches(start_x)
    y = start_y
    shapes = []
    for step in steps:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(y), Inches(box_w), Inches(box_h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xE0,0xF2,0xFF)
        shp.line.color.rgb = RGBColor(0x00,0x66,0x99)
        tx = shp.text_frame
        tx.text = step
        tx.paragraphs[0].font.size = Pt(14)
        tx.paragraphs[0].font.bold = True
        shapes.append(shp)
        y += box_h + gap
    # add connectors
    for i in range(len(shapes)-1):
        start = shapes[i].bottom_center
        end = shapes[i+1].top_center
        # Add a simple line (connector) by adding a shape of type line
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start.x, start.y, end.x, end.y)
    return slide

# Title
add_title_slide("Transient Stability Assessment Using BCU + LSTM", "Based on Bhui & Senroy (2016) — Improved with LSTM prediction")

# Motivation
add_bullets_slide("Motivation", [
    "Transient instability risks: generator loss of synchronism and blackouts.",
    "TEF + BCU provides a fast analytical stability margin.",
    "Bhui & Senroy used polynomial curve fitting for short-term state prediction.",
    "Our contribution: Replace curve-fitting with LSTM for better prediction accuracy."
])

# Architecture
add_bullets_slide("Overall Architecture", [
    "Two stages: Offline (database creation) and Online (real-time assessment)",
    "Offline: PSSE simulations → BCU → TEF calculations → Lookup database",
    "Online: PMU/measurements → KE_norm matching → TEF margin → LSTM prediction → Control action",
    "Final system: IEEE-39 bus (IEEE-9 for development)"
])

# Offline Stage slide
add_bullets_slide("Offline Stage Details", [
    "Simulate faults with PSSE (various locations, durations, clearings).",
    "Extract δ(t), ω(t), Pe(t) in COI frame for each generator.",
    "Apply BCU to find MOD, SEP, and CUEP.",
    "Compute TEF quantities: KE_norm, ΔV_PE, V_cr, CCT.",
    "Store results into a lookup database for online matching."
])

# Offline flowchart (graphical)
offline_steps = [
    "PSSE: Multi-Fault Simulation",
    "Convert trajectories to COI frame",
    "BCU: Identify MOD / SEP / CUEP",
    "Compute TEF metrics (KE_norm, ΔV_PE, V_cr)",
    "Build Lookup Database (MOD + KE_norm)"
]
add_flowchart_slide("Offline Stage Flowchart", offline_steps, start_x=1.2, start_y=1.6, box_w=10, box_h=0.6, gap=0.4)

# TEF outputs
add_bullets_slide("TEF & Database Outputs", [
    "Stored per fault location / topology:",
    "- Mode of Disturbance (MOD)",
    "- Stable and Unstable Equilibrium Points (SEP, CUEP)",
    "- KE_norm vectors, ΔV_PE curves, V_cr, CCT",
    "Lookup DB used for fast online matching"
])

# Online Stage slide
add_bullets_slide("Online Stage Details", [
    "Acquire δ(t), ω(t) from PMUs or PSSE channels (COI-referenced).",
    "Compute normalized KE immediately after fault clearing.",
    "Find probable MOD by matching KE_norm vector against lookup DB (k-NN).",
    "Compute TEF stability margin for candidate MODs.",
    "Predict δ(t+τ) and ω(t+τ) using LSTM (our contribution).",
    "Decide control action (generation shedding) using predicted states + TEF."
])

# Online flowchart (graphical)
online_steps = [
    "Real-time Measurements (δ, ω)",
    "Convert to COI Frame",
    "Compute KE_norm",
    "MOD Matching (lookup DB)",
    "Compute TEF Margin",
    "LSTM Prediction of δ(t+τ), ω(t+τ)",
    "Control Action Determination (shedding)"
]
add_flowchart_slide("Online Stage Flowchart", online_steps, start_x=1.2, start_y=1.2, box_w=10, box_h=0.55, gap=0.35)

# Contribution slide with diagram idea
add_bullets_slide("Our Contribution: LSTM Integration", [
    "Replaces quadratic/cubic curve fitting used in base paper.",
    "Inputs to LSTM: short window of δ(t), ω(t), Pe(t) (0–0.1s).",
    "Output: δ(t+τ), ω(t+τ) for τ ≈ 0.2–0.3s (prediction horizon).",
    "Benefits: models nonlinearity, improves prediction error, reduces incorrect control actions."
])

# LSTM details slide
add_bullets_slide("LSTM Model & Data Pipeline", [
    "Training Data: Offline PSSE simulations (many fault cases).",
    "Feature engineering: COI-corrected angle & speed, normalized KE, windowing.",
    "Network: small stacked LSTM (1–2 layers) with dropout, trained in MATLAB/Python.",
    "Evaluation: compare predicted δ vs ground truth; target < 2° error."
])

# Tools slide
add_bullets_slide("Tools & Implementation", [
    "PSSE 36.x for dynamic simulations and channels (δ, ω, Pe).",
    "MATLAB for TEF calculations, BCU routines, LSTM training (or Python).",
    "Python (psspy) for automation and batch simulations.",
    "Data exchange: CSV/Excel → MATLAB/Python for ML pipeline."
])

# Expected contributions slide
add_bullets_slide("Expected Contributions & Results", [
    "Improved rotor angle prediction accuracy over polynomial fits.",
    "More reliable TEF-based stability margins.",
    "Reduced unnecessary/insufficient generation shedding.",
    "Demonstration on IEEE-39 bus with quantitative error & performance metrics."
])

# Conclusion slide
add_bullets_slide("Conclusion", [
    "Methodology matches Bhui & Senroy's framework (BCU + TEF).",
    "LSTM improves the single missing block (prediction module).",
    "PSSE provides industry-grade dynamic fidelity.",
    "Ready for implementation: offline DB → online LSTM → control action."
])

# Q/A slide
add_bullets_slide("Q/A", [
    "Thank you — Questions?"
])

output_path = "/mnt/data/Graphical_Methodology_Presentation.pptx"
prs.save(output_path)
output_path
